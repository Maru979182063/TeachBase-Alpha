import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path

THIS_DIR = Path(__file__).resolve().parent
VENDOR_DIR = THIS_DIR / "vendor"
if VENDOR_DIR.exists():
    sys.path.insert(0, str(VENDOR_DIR))

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.pdfmetrics import registerFont
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


VERSION_ORDER = ["基础版", "常用版", "进阶版"]
AUDIENCE_ORDER = ["学生版", "教师版"]


PUA_RE = re.compile(r"[\ue000-\uf8ff]")


def safe_name(text: str) -> str:
    text = str(text or "").strip()
    text = re.sub(r'[\\/:*?"<>|]+', "_", text)
    text = re.sub(r"\s+", "_", text)
    return text or "未命名"


def set_doc_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_default_style(document: Document) -> None:
    normal = document.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.72)
    section.right_margin = Inches(0.72)


def normalize_preview_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def preview_looks_noisy(text: str) -> bool:
    normalized = normalize_preview_text(text)
    if not normalized:
        return True

    cjk_count = len(re.findall(r"[\u4e00-\u9fff]", normalized))
    latin_count = len(re.findall(r"[A-Za-z]", normalized))
    digit_count = len(re.findall(r"\d", normalized))
    math_symbol_count = len(re.findall(r"[=+\-*/^<>]", normalized))
    private_use_count = len(PUA_RE.findall(normalized))
    symbol_count = len(normalized) - cjk_count - latin_count - digit_count
    spaced_operator_runs = len(
        re.findall(r"(?:[=+\-*/^(){}\[\]<>_.,:;|\\]+(?:\s+|$)){6,}", normalized)
    )
    sparse_readable = cjk_count <= 8 and latin_count + digit_count + symbol_count > 40
    symbol_heavy = len(normalized) >= 48 and symbol_count / max(len(normalized), 1) > 0.42 and cjk_count < 16
    formula_noise = len(normalized) >= 72 and math_symbol_count >= 10 and cjk_count <= 36
    return private_use_count >= 1 or spaced_operator_runs >= 1 or sparse_readable or symbol_heavy or formula_noise


def export_preview_text(question: dict) -> str:
    preview = normalize_preview_text(question.get("previewText", ""))
    storage_mode = str(question.get("textStorageMode") or "")
    if preview and storage_mode != "ocr_reference_only" and not preview_looks_noisy(preview):
        return preview

    checkpoint = str(question.get("checkpoint") or "\u5f53\u524d\u9898\u5757")
    component = str(question.get("componentLabel") or "\u9898\u56fe\u9884\u89c8")
    source_page = question.get("sourcePage")
    source_page_label = f"P{source_page}" if source_page else "\u9875\u7801\u5f85\u590d\u6838"
    return (
        f"{checkpoint}\uff5c{component}\uff5c{source_page_label}"
        "\u3002\u5f53\u524d\u4ee5\u9898\u56fe\u4e3a\u51c6\uff0c\u6587\u5b57\u5c42\u5f85\u590d\u6838\u3002"
    )


def make_docx(payload: dict, target_path: Path) -> None:
    lesson = payload["lesson"]
    questions = payload["questions"]
    audience = payload["audience"]
    version = payload["version"]
    document = Document()
    set_default_style(document)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(f"{lesson['lesson_title']} · {version} · {audience}")
    run.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x1F, 0x59, 0xD7)

    meta = document.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_run = meta.add_run(
        f"{lesson['stage']}数学 / {lesson['grade']} / {lesson['season']} / 第{lesson['lesson_no']}讲"
    )
    meta_run.font.size = Pt(10.5)
    meta_run.font.color.rgb = RGBColor(0x6D, 0x7A, 0x8C)

    summary_table = document.add_table(rows=2, cols=4)
    summary_table.style = "Table Grid"
    summary_table.autofit = True
    headers = ["题量", "知识点数", "来源讲义", "导出时间"]
    values = [
        str(len(questions)),
        str(lesson.get("knowledge_point_count", 0)),
        lesson.get("source_pdf_name", ""),
        payload["created_at_display"],
    ]
    for idx, header in enumerate(headers):
        summary_table.cell(0, idx).text = header
        set_doc_cell_shading(summary_table.cell(0, idx), "EDF4FF")
        summary_table.cell(1, idx).text = values[idx]

    document.add_paragraph()
    heading = document.add_paragraph()
    heading_run = heading.add_run("课程目标")
    heading_run.bold = True
    heading_run.font.size = Pt(13)
    for line in [line.strip() for line in lesson.get("objectives", "").splitlines() if line.strip()]:
        document.add_paragraph(line, style="List Bullet")

    document.add_paragraph()
    heading = document.add_paragraph()
    heading_run = heading.add_run("知识树概览")
    heading_run.bold = True
    heading_run.font.size = Pt(13)
    for module in payload.get("knowledge_tree", []):
        module_p = document.add_paragraph(style="List Bullet")
        module_p.add_run(module.get("module", "未命名模块")).bold = True
        for item in module.get("items", []):
            document.add_paragraph(str(item), style="List Bullet 2")

    document.add_section(WD_SECTION.NEW_PAGE)
    q_title = document.add_paragraph()
    q_run = q_title.add_run("题块清单")
    q_run.bold = True
    q_run.font.size = Pt(15)

    for index, question in enumerate(questions, start=1):
        p = document.add_paragraph()
        p_format = p.paragraph_format
        p_format.space_before = Pt(10)
        p_format.space_after = Pt(6)
        r = p.add_run(f"{index}. {question['localNumber']}｜{question['checkpoint']}")
        r.bold = True
        r.font.size = Pt(12.5)
        meta_p = document.add_paragraph()
        meta_p.add_run(
            f"{question['componentLabel']} · P{question['sourcePage']} · 标签：{' / '.join(question['effectiveVersionTags'])}"
        ).font.color.rgb = RGBColor(0x6D, 0x7A, 0x8C)

        crop_path = question.get("cropPath")
        if crop_path and os.path.exists(crop_path):
            document.add_picture(crop_path, width=Inches(5.85))

        if audience == "教师版":
            prompt = document.add_paragraph()
            prompt.add_run("题块摘要：").bold = True
            prompt.add_run(export_preview_text(question))
            review = document.add_paragraph()
            review.add_run("视觉备注：").bold = True
            review.add_run(question.get("reviewNote", ""))
            if question.get("riskIssues"):
                risk = document.add_paragraph()
                risk.add_run("风险提示：").bold = True
                risk.add_run("；".join(question["riskIssues"]))
        else:
            student_note = document.add_paragraph()
            student_note.add_run("保留题块原始视觉切片，便于继续组装学生版讲义。").italic = True

    document.save(str(target_path))


def make_pdf(payload: dict, target_path: Path) -> None:
    lesson = payload["lesson"]
    questions = payload["questions"]
    audience = payload["audience"]
    version = payload["version"]

    registerFont(UnicodeCIDFont("STSong-Light"))
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleCN",
        parent=styles["Title"],
        fontName="STSong-Light",
        fontSize=20,
        leading=26,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#1f59d7"),
    )
    body_style = ParagraphStyle(
        "BodyCN",
        parent=styles["BodyText"],
        fontName="STSong-Light",
        fontSize=10.5,
        leading=16,
        textColor=colors.HexColor("#1d2736"),
    )
    meta_style = ParagraphStyle(
        "MetaCN",
        parent=styles["BodyText"],
        fontName="STSong-Light",
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor("#6d7a8c"),
    )
    heading_style = ParagraphStyle(
        "HeadingCN",
        parent=styles["Heading2"],
        fontName="STSong-Light",
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#1d2736"),
    )

    story = [
        Paragraph(f"{lesson['lesson_title']} · {version} · {audience}", title_style),
        Spacer(1, 6 * mm),
        Paragraph(
            f"{lesson['stage']}数学 / {lesson['grade']} / {lesson['season']} / 第{lesson['lesson_no']}讲 / 导出时间 {payload['created_at_display']}",
            meta_style,
        ),
        Spacer(1, 6 * mm),
    ]

    table_data = [
        ["题量", "知识点数", "来源讲义", "导出方式"],
        [str(len(questions)), str(lesson.get("knowledge_point_count", 0)), lesson.get("source_pdf_name", ""), "题块拆分导出"],
    ]
    table = Table(table_data, colWidths=[24 * mm, 24 * mm, 95 * mm, 35 * mm])
    table.setStyle(
        TableStyle(
            [
                ("FONTNAME", (0, 0), (-1, -1), "STSong-Light"),
                ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#edf4ff")),
                ("GRID", (0, 0), (-1, -1), 0.6, colors.HexColor("#d8e4f2")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]
        )
    )
    story.extend([table, Spacer(1, 7 * mm), Paragraph("课程目标", heading_style)])

    for line in [line.strip() for line in lesson.get("objectives", "").splitlines() if line.strip()]:
        story.append(Paragraph(f"• {line}", body_style))

    story.extend([Spacer(1, 5 * mm), Paragraph("知识树概览", heading_style)])
    for module in payload.get("knowledge_tree", []):
        story.append(Paragraph(f"• {module.get('module', '未命名模块')}", body_style))
        for item in module.get("items", []):
            story.append(Paragraph(f"　- {item}", meta_style))

    story.append(PageBreak())
    story.append(Paragraph("题块清单", heading_style))

    for index, question in enumerate(questions, start=1):
        story.append(Spacer(1, 4 * mm))
        story.append(Paragraph(f"{index}. {question['localNumber']}｜{question['checkpoint']}", heading_style))
        story.append(
            Paragraph(
                f"{question['componentLabel']} · P{question['sourcePage']} · 标签：{' / '.join(question['effectiveVersionTags'])}",
                meta_style,
            )
        )
        crop_path = question.get("cropPath")
        if crop_path and os.path.exists(crop_path):
            image = Image(crop_path)
            image._restrictSize(172 * mm, 95 * mm)
            story.append(Spacer(1, 2 * mm))
            story.append(image)
        if audience == "教师版":
            story.append(Spacer(1, 2 * mm))
            story.append(Paragraph(f"题块摘要：{export_preview_text(question)}", body_style))
            story.append(Paragraph(f"视觉备注：{question.get('reviewNote', '')}", meta_style))
            if question.get("riskIssues"):
                story.append(Paragraph(f"风险提示：{'；'.join(question['riskIssues'])}", meta_style))
        story.append(Spacer(1, 4 * mm))

    doc = SimpleDocTemplate(
        str(target_path),
        pagesize=A4,
        leftMargin=15 * mm,
        rightMargin=15 * mm,
        topMargin=14 * mm,
        bottomMargin=12 * mm,
    )
    doc.build(story)


def build_variant_payload(base_payload: dict, version: str, audience: str) -> dict:
    lesson = base_payload["lesson"]
    questions = [
        q
        for q in base_payload["splitLesson"]["questions"]
        if version in q.get("effectiveVersionTags", q.get("versionTags", []))
    ]
    return {
        "lesson": lesson,
        "version": version,
        "audience": audience,
        "questions": questions,
        "knowledge_tree": base_payload["splitLesson"].get("tree", []),
        "created_at_display": base_payload["createdAtDisplay"],
    }


def add_ppt_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, bold: bool = False, color: str = "1D2736", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)
    return box


def add_ppt_card(slide, left: float, top: float, width: float, height: float, title: str, value: str, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor.from_string("FFFFFF")
    shape.line.color.rgb = PptRGBColor.from_string("D8E4F2")
    add_ppt_text(slide, title, left + 0.14, top + 0.1, width - 0.2, 0.22, 10, False, "6D7A8C")
    add_ppt_text(slide, str(value), left + 0.14, top + 0.42, width - 0.2, 0.3, 22, True, color)


def make_compass_ppt(payload: dict, target_path: Path) -> None:
    lesson = payload["lesson"]
    split_lesson = payload["splitLesson"]
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_title(slide, title: str, subtitle: str = ""):
        add_ppt_text(slide, title, 0.6, 0.38, 8.8, 0.42, 24, True)
        if subtitle:
            add_ppt_text(slide, subtitle, 0.62, 0.86, 10.2, 0.25, 10, False, "6D7A8C")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "讲义拆分导出罗盘", f"{lesson['grade']} · {lesson['season']} · {lesson['lesson_title']}")
    add_ppt_text(
        slide,
        "本次导出聚焦于“按课次 + 最小知识点”的讲义重组，方便老师沿着原知识体系继续拆分与组装。",
        0.72,
        1.45,
        11.2,
        0.8,
        18,
        False,
        "1D2736",
    )
    add_ppt_card(slide, 0.72, 3.05, 2.5, 1.0, "总题块数", str(split_lesson.get("question_count", len(split_lesson["questions"]))), "2D6DF6")
    add_ppt_card(slide, 3.48, 3.05, 2.5, 1.0, "知识点数", str(lesson.get("knowledge_point_count", 0)), "1EA76A")
    add_ppt_card(slide, 6.24, 3.05, 2.5, 1.0, "视觉已审", str(split_lesson.get("auditSummary", {}).get("reviewedCount", 0)), "F59B23")
    add_ppt_card(slide, 9.0, 3.05, 2.5, 1.0, "待人工关注", str(split_lesson.get("auditSummary", {}).get("pendingCount", 0)), "F05555")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "课次底盘", "把讲义身份、来源与课程目标收束成一页，适合老师快速对焦。")
    add_ppt_text(slide, f"学段：{lesson['stage']}\n年级：{lesson['grade']}\n季节：{lesson['season']}\n讲次：第{lesson['lesson_no']}讲\n来源：{lesson['source_pdf_name']}", 0.8, 1.45, 3.5, 3.2, 16)
    add_ppt_text(slide, "课程目标", 4.65, 1.45, 2.0, 0.3, 17, True)
    add_ppt_text(slide, "\n".join([f"• {line.strip()}" for line in lesson.get("objectives", "").splitlines() if line.strip()]), 4.7, 1.95, 7.5, 3.8, 15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "知识树与版本分布", "这一页专门展示当前课次的结构根茎，以及三个版本的题量差异。")
    tree_lines = []
    for module in split_lesson.get("tree", []):
        tree_lines.append(f"• {module.get('module', '未命名模块')}")
        for item in module.get("items", []):
            tree_lines.append(f"  - {item}")
    add_ppt_text(slide, "知识树摘要", 0.82, 1.38, 2.5, 0.25, 16, True)
    add_ppt_text(slide, "\n".join(tree_lines[:18]), 0.86, 1.8, 6.2, 4.9, 13)

    stats = {"基础版": 0, "常用版": 0, "进阶版": 0}
    for question in split_lesson["questions"]:
        for tag in question.get("effectiveVersionTags", question.get("versionTags", [])):
            if tag in stats:
                stats[tag] += 1
    add_ppt_text(slide, "版本题量", 8.05, 1.38, 2.0, 0.25, 16, True)
    add_ppt_card(slide, 8.08, 1.88, 3.5, 1.0, "基础版", str(stats["基础版"]), "2D6DF6")
    add_ppt_card(slide, 8.08, 3.12, 3.5, 1.0, "常用版", str(stats["常用版"]), "F59B23")
    add_ppt_card(slide, 8.08, 4.36, 3.5, 1.0, "进阶版", str(stats["进阶版"]), "1EA76A")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "风险关注点", "把需要人工盯一下的题块提前捞出来，演示时会更直观。")
    risky = payload.get("reviewItems", [])[:4]
    if not risky:
        risky = [q for q in split_lesson["questions"] if q.get("risk") != "低风险"][:4]
    pos = [(0.72, 1.5), (6.72, 1.5), (0.72, 4.0), (6.72, 4.0)]
    for idx, item in enumerate(risky):
        left, top = pos[idx]
        add_ppt_text(
            slide,
            f"{item.get('localNumber', item.get('queueNo', ''))}｜{item.get('checkpoint', item.get('title', ''))}",
            left,
            top,
            4.9,
            0.3,
            12,
            True,
        )
        crop_path = item.get("cropPath")
        if crop_path and Path(crop_path).exists():
            slide.shapes.add_picture(str(crop_path), PptInches(left), PptInches(top + 0.4), width=PptInches(2.2))
        note = item.get("reviewNote") or "；".join(item.get("tags", [])) or "建议人工再看一下视觉边界。"
        add_ppt_text(slide, note[:90], left + 2.35, top + 0.46, 2.8, 1.2, 10, False, "1D2736")
        add_ppt_text(slide, item.get("risk", "待审"), left + 4.4, top, 0.7, 0.2, 11, True, "F05555" if item.get("risk") == "高风险" else "F59B23", PP_ALIGN.RIGHT)

    prs.save(str(target_path))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--payload", required=True)
    parser.add_argument("--manifest", required=True)
    args = parser.parse_args()

    payload = json.loads(Path(args.payload).read_text(encoding="utf-8"))
    output_dir = Path(payload["outputDir"])
    output_dir.mkdir(parents=True, exist_ok=True)

    manifest_files = []
    for version in [v for v in VERSION_ORDER if v in payload["selectedVersions"]]:
        for audience in [a for a in AUDIENCE_ORDER if a in payload["selectedAudiences"]]:
            item_payload = build_variant_payload(payload, version, audience)
            base_name = safe_name(
                f"{payload['lesson']['stage']}数学_{payload['lesson']['grade']}_{payload['lesson']['season']}_{payload['lesson']['lesson_title']}_{version}_{audience}"
            )
            for fmt in payload["selectedFormats"]:
                target_path = output_dir / f"{base_name}.{fmt.lower()}"
                if fmt == "DOCX":
                    make_docx(item_payload, target_path)
                elif fmt == "PDF":
                    make_pdf(item_payload, target_path)
                else:
                    continue
                manifest_files.append(
                    {
                        "name": target_path.name,
                        "path": str(target_path),
                        "format": fmt,
                        "version": version,
                        "audience": audience,
                        "questionCount": len(item_payload["questions"]),
                    }
                )

    if payload.get("includeCompass"):
        compass_name = safe_name(
            f"{payload['lesson']['stage']}数学_{payload['lesson']['grade']}_{payload['lesson']['season']}_{payload['lesson']['lesson_title']}_讲义罗盘"
        )
        compass_path = output_dir / f"{compass_name}.pptx"
        make_compass_ppt(payload, compass_path)
        manifest_files.append(
            {
                "name": compass_path.name,
                "path": str(compass_path),
                "format": "PPTX",
                "version": "讲义罗盘",
                "audience": "工作台",
                "questionCount": len(payload["splitLesson"]["questions"]),
            }
        )

    Path(args.manifest).write_text(
        json.dumps({"files": manifest_files, "generatedAt": datetime.now().isoformat()}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
